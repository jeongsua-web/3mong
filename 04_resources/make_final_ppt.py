# -*- coding: utf-8 -*-
"""Fluento 최종발표 PPT 초안 생성 스크립트"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ---------- 팔레트 ----------
PRIMARY = RGBColor(0x4F, 0x46, 0xE5)   # indigo
DARK    = RGBColor(0x1E, 0x1B, 0x4B)   # deep navy
ACCENT  = RGBColor(0x8B, 0x5C, 0xF6)   # violet
LIGHT   = RGBColor(0xF5, 0xF4, 0xFF)   # light bg
GRAY    = RGBColor(0x6B, 0x72, 0x80)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
GREEN   = RGBColor(0x10, 0xB9, 0x81)
RED     = RGBColor(0xEF, 0x44, 0x44)
AMBER   = RGBColor(0xF5, 0x9E, 0x0B)

FONT = "Apple SD Gothic Neo"

SW, SH = Inches(13.333), Inches(7.5)

prs = Presentation()
prs.slide_width = SW
prs.slide_height = SH
BLANK = prs.slide_layouts[6]


def set_font(run, size, color=DARK, bold=False, name=FONT):
    f = run.font
    f.size = Pt(size)
    f.color.rgb = color
    f.bold = bold
    f.name = name
    # 한글 폰트 지정
    rPr = run._r.get_or_add_rPr()
    ea = rPr.find(qn('a:ea'))
    if ea is None:
        ea = rPr.makeelement(qn('a:ea'), {})
        rPr.append(ea)
    ea.set('typeface', name)


def add_box(slide, x, y, w, h, fill=None, line=None, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.08, shadow=False):
    sp = slide.shapes.add_shape(shape, x, y, w, h)
    if shape == MSO_SHAPE.ROUNDED_RECTANGLE:
        try:
            sp.adjustments[0] = radius
        except Exception:
            pass
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid()
        sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = Pt(1)
    sp.shadow.inherit = False
    return sp


def add_text(slide, x, y, w, h, lines, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    """lines: list of (text, size, color, bold) or list of runs-lists"""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        runs = line if isinstance(line, list) else [line]
        for (text, size, color, bold) in runs:
            r = p.add_run()
            r.text = text
            set_font(r, size, color, bold)
        p.space_after = Pt(6)
    return tb


def content_slide(label, title):
    """공통 레이아웃: 상단 라벨 + 제목 + 구분선"""
    s = prs.slides.add_slide(BLANK)
    add_box(s, 0, 0, SW, SH, fill=WHITE, shape=MSO_SHAPE.RECTANGLE)
    add_box(s, 0, 0, Inches(0.18), SH, fill=PRIMARY, shape=MSO_SHAPE.RECTANGLE)
    add_text(s, Inches(0.7), Inches(0.45), Inches(11), Inches(0.4),
             [(label, 13, ACCENT, True)])
    add_text(s, Inches(0.7), Inches(0.8), Inches(12), Inches(0.8),
             [(title, 30, DARK, True)])
    add_box(s, Inches(0.72), Inches(1.55), Inches(2.2), Pt(3), fill=PRIMARY, shape=MSO_SHAPE.RECTANGLE)
    return s


def bullets(slide, x, y, w, items, size=16, gap=10, body_color=DARK):
    """items: list of (head, body) — head는 굵게, body는 일반"""
    tb = slide.shapes.add_textbox(x, y, w, Inches(5))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    first = True
    for head, body in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        r = p.add_run(); r.text = "▪  "; set_font(r, size, PRIMARY, True)
        if head:
            r = p.add_run(); r.text = head; set_font(r, size, DARK, True)
        if body:
            r = p.add_run(); r.text = ("  " if head else "") + body
            set_font(r, size, body_color if head else DARK, False)
        p.space_after = Pt(gap)
    return tb


def card(slide, x, y, w, h, title, desc, fill=LIGHT, title_color=PRIMARY, title_size=15, desc_size=12.5):
    add_box(slide, x, y, w, h, fill=fill)
    add_text(slide, x + Inches(0.22), y + Inches(0.16), w - Inches(0.44), Inches(0.5),
             [(title, title_size, title_color, True)])
    add_text(slide, x + Inches(0.22), y + Inches(0.62), w - Inches(0.44), h - Inches(0.75),
             [(desc, desc_size, GRAY, False)])


def arrow(slide, x, y, w, h=Pt(20)):
    a = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x, y, w, h)
    a.fill.solid(); a.fill.fore_color.rgb = ACCENT
    a.line.fill.background()
    a.shadow.inherit = False
    return a


def diagram_node(slide, x, y, w, h, title, sub, fill=PRIMARY, text_color=WHITE):
    add_box(slide, x, y, w, h, fill=fill, radius=0.15)
    add_text(slide, x, y + Inches(0.08), w, Inches(0.4),
             [(title, 14, text_color, True)], align=PP_ALIGN.CENTER)
    if sub:
        add_text(slide, x, y + Inches(0.46), w, h - Inches(0.5),
                 [(sub, 10.5, text_color, False)], align=PP_ALIGN.CENTER)


# ============================================================
# 1. 표지
# ============================================================
s = prs.slides.add_slide(BLANK)
add_box(s, 0, 0, SW, SH, fill=DARK, shape=MSO_SHAPE.RECTANGLE)
add_box(s, 0, Inches(7.1), SW, Inches(0.4), fill=PRIMARY, shape=MSO_SHAPE.RECTANGLE)
add_box(s, Inches(9.8), Inches(-1.2), Inches(5), Inches(5), fill=PRIMARY, shape=MSO_SHAPE.OVAL)
add_box(s, Inches(10.8), Inches(-0.5), Inches(4), Inches(4), fill=ACCENT, shape=MSO_SHAPE.OVAL)
add_text(s, Inches(1), Inches(2.1), Inches(10), Inches(0.5),
         [("최종 발표", 16, ACCENT, True)])
add_text(s, Inches(1), Inches(2.6), Inches(11), Inches(1.4),
         [("Fluento", 60, WHITE, True)])
add_text(s, Inches(1), Inches(3.9), Inches(11), Inches(0.7),
         [("AI 친구와 대화하며 배우는 영어 회화 학습 서비스", 22, WHITE, False)])
add_text(s, Inches(1), Inches(5.6), Inches(10), Inches(1),
         [("팀명  ·  발표자 이름", 14, RGBColor(0xC7, 0xC5, 0xE8), False),
          ("2026. 06.", 14, RGBColor(0xC7, 0xC5, 0xE8), False)])

# ============================================================
# 2. 목차
# ============================================================
s = content_slide("AGENDA", "발표 순서")
items = [
    ("01", "기획 배경", "왜 AI 영어 회화인가"),
    ("02", "핵심 기능", "Fluento가 제공하는 학습 경험"),
    ("03", "시스템 아키텍처 & 기술 스택", "AWS Bedrock · Cognito · EC2"),
    ("04", "구현 상세", "AI 대화 설계 · 레벨 평가 · 인증"),
    ("05", "품질 검증 & 트러블슈팅", "E2E 72건 테스트와 버그 수정기"),
    ("06", "라이브 시연", "회원가입부터 오답노트까지"),
    ("07", "회고 & 향후 계획", ""),
]
y = Inches(2.0)
for num, t, d in items:
    add_text(s, Inches(1.0), y, Inches(0.8), Inches(0.5), [(num, 18, ACCENT, True)])
    add_text(s, Inches(1.9), y, Inches(4.5), Inches(0.5), [(t, 18, DARK, True)])
    add_text(s, Inches(6.6), y + Inches(0.04), Inches(6), Inches(0.5), [(d, 13, GRAY, False)])
    y += Inches(0.72)

# ============================================================
# 3. 기획 배경
# ============================================================
s = content_slide("01  BACKGROUND", "왜 만들었나 — 영어 회화의 진짜 장벽")
bullets(s, Inches(0.9), Inches(2.0), Inches(11.5), [
    ("말할 상대가 없다.", "문법·단어 공부는 혼자 가능하지만, 회화는 \"실제로 말해볼 상대\"가 필요"),
    ("기존 서비스의 한계.", "전화영어는 비싸고 예약이 필요, 단어 앱은 회화로 이어지지 않음"),
    ("틀려도 부끄럽지 않은 상대.", "사람 앞에서는 실수가 두렵지만 AI 앞에서는 부담 없이 시도 가능"),
], size=17, gap=14)
card(s, Inches(0.9), Inches(4.6), Inches(11.5), Inches(2.0),
     "Fluento의 해법",
     "AWS Bedrock(Claude) 기반 AI 친구와 24시간 영어로 대화하고, 대화 중 실수는 자동으로 오답노트에 쌓이며, "
     "실력에 맞춰 대화 난이도가 조절되는 개인 맞춤 영어 회화 서비스",
     fill=LIGHT, title_size=16, desc_size=14)

# ============================================================
# 4. 핵심 기능
# ============================================================
s = content_slide("02  FEATURES", "핵심 기능 한눈에")
feats = [
    ("AI 친구와 영어 채팅", "Claude 3.5 Sonnet 기반, SSE 스트리밍으로 실시간 응답"),
    ("레벨 자동 평가·조절", "대화 내용을 분석해 사용자 실력에 맞게 난이도 자동 조정"),
    ("오답노트", "대화 중 틀린 표현을 AI가 감지해 자동 수집 → 복습"),
    ("오늘의 목표 & 통계", "학습 시간·암기 단어·연속 학습일 자동 추적"),
    ("커스텀 AI 친구", "성격·말투를 직접 설정한 나만의 대화 상대 생성"),
    ("소셜 로그인", "Google 계정으로 간편하게 시작 (Cognito 페더레이션)"),
]
xs = [Inches(0.8), Inches(5.0), Inches(9.2)]
ys = [Inches(2.0), Inches(4.45)]
for i, (t, d) in enumerate(feats):
    card(s, xs[i % 3], ys[i // 3], Inches(3.9), Inches(2.1), t, d)

# ============================================================
# 5. 시스템 아키텍처
# ============================================================
s = content_slide("03  ARCHITECTURE", "시스템 아키텍처")
# 행 1: 클라이언트 → 프론트 → 백엔드 → Bedrock
ny = Inches(2.3); nh = Inches(1.1)
diagram_node(s, Inches(0.8),  ny, Inches(2.1), nh, "사용자", "Web Browser", fill=GRAY)
arrow(s, Inches(3.0), ny + Inches(0.45), Inches(0.55))
diagram_node(s, Inches(3.65), ny, Inches(2.3), nh, "Frontend", "React + Vite", fill=ACCENT)
arrow(s, Inches(6.05), ny + Inches(0.45), Inches(0.55))
diagram_node(s, Inches(6.7),  ny, Inches(2.6), nh, "Backend", "Spring Boot 4 · WebFlux SSE", fill=PRIMARY)
arrow(s, Inches(9.4), ny + Inches(0.45), Inches(0.55))
diagram_node(s, Inches(10.05), ny, Inches(2.5), nh, "AWS Bedrock", "Claude 3.5 Sonnet", fill=DARK)
# 행 2: 백엔드 하단 의존성
ny2 = Inches(4.1)
diagram_node(s, Inches(3.65), ny2, Inches(2.3), nh, "AWS Cognito", "JWT 인증 · Google IdP", fill=LIGHT, text_color=DARK)
diagram_node(s, Inches(6.2),  ny2, Inches(2.0), nh, "PostgreSQL 17", "사용자·대화·오답", fill=LIGHT, text_color=DARK)
diagram_node(s, Inches(8.45), ny2, Inches(1.9), nh, "Redis 7", "캐시·세션", fill=LIGHT, text_color=DARK)
diagram_node(s, Inches(10.55), ny2, Inches(2.0), nh, "EC2 + Docker", "GitHub Actions 배포", fill=LIGHT, text_color=DARK)
add_text(s, Inches(0.9), Inches(5.7), Inches(11.6), Inches(1.3), [
    [("핵심 의사결정  ", 14, PRIMARY, True),
     ("① SSE 스트리밍으로 AI 응답을 토큰 단위 전달 → 체감 대기시간 최소화   "
      "② Bedrock 리전 분리로 모델 가용성 확보   ③ 인증은 Cognito에 위임해 보안 구현 비용 절감", 13.5, DARK, False)],
])

# ============================================================
# 6. 기술 스택
# ============================================================
s = content_slide("03  TECH STACK", "기술 스택")
rows = [
    ("Frontend", "React 18 · Vite · Playwright(E2E)"),
    ("Backend", "Java 17 · Spring Boot 4 · Spring Security 7 · WebFlux(SSE) · JPA · Flyway"),
    ("AI", "AWS Bedrock — Claude 3.5 Sonnet"),
    ("Auth", "AWS Cognito (JWT, OAuth2 Resource Server) · Google 소셜 로그인"),
    ("Database", "PostgreSQL 17 · Redis 7"),
    ("Infra / CI·CD", "AWS EC2 · Docker Compose · GitHub Actions 자동 배포"),
    ("Docs / 협업", "Swagger UI (springdoc-openapi) · GitHub Monorepo"),
]
y = Inches(2.0)
for k, v in rows:
    add_box(s, Inches(0.9), y, Inches(2.6), Inches(0.55), fill=PRIMARY, radius=0.2)
    add_text(s, Inches(0.9), y + Inches(0.1), Inches(2.6), Inches(0.4),
             [(k, 14, WHITE, True)], align=PP_ALIGN.CENTER)
    add_text(s, Inches(3.8), y + Inches(0.1), Inches(8.8), Inches(0.4),
             [(v, 14, DARK, False)])
    y += Inches(0.68)

# ============================================================
# 7. AI 대화 설계
# ============================================================
s = content_slide("04  AI DESIGN", "AI 대화 설계 — 페르소나와 난이도 제어")
bullets(s, Inches(0.9), Inches(2.0), Inches(11.6), [
    ("AI 친구 페르소나 프롬프트.", "캐릭터별 성격·말투·관심사를 시스템 프롬프트로 정의 → 친구마다 다른 대화 경험"),
    ("레벨별 응답 난이도 제어.", "사용자 레벨(어휘·문장 길이·표현 복잡도)을 프롬프트에 주입해 응답 난이도를 동적으로 조절"),
    ("교정 피드백 내장.", "사용자의 문법·표현 오류를 감지하면 자연스럽게 교정하고, 오류 데이터는 오답노트로 적재"),
    ("SSE 스트리밍.", "Bedrock 응답을 WebFlux SSE로 토큰 단위 중계 → 첫 글자까지의 대기시간 단축"),
], size=15.5, gap=12)
card(s, Inches(0.9), Inches(5.1), Inches(11.6), Inches(1.6),
     "대화 → 평가 → 반영 파이프라인",
     "대화 종료(또는 일정 메시지 누적) 시 비동기로 레벨 평가 요청 → 평가 점수를 사용자 레벨에 반영 → "
     "다음 대화부터 조정된 난이도 적용",
     fill=LIGHT, title_size=15, desc_size=13.5)

# ============================================================
# 8. 인증 & 보안
# ============================================================
s = content_slide("04  AUTH", "인증 & 보안 — AWS Cognito")
bullets(s, Inches(0.9), Inches(2.0), Inches(11.6), [
    ("JWT 기반 인증.", "Spring Security OAuth2 Resource Server가 Cognito 발급 JWT를 검증"),
    ("이메일 인증 회원가입.", "표준 Cognito signUp → confirmSignUp 흐름으로 이메일 검증"),
    ("Google 소셜 로그인.", "Cognito Hosted UI에 Google IdP 페더레이션 — 백엔드 코드 변경 0줄로 연동"),
    ("인증 책임 분리.", "비밀번호 저장·토큰 발급·소셜 연동을 모두 Cognito에 위임 → 직접 구현 대비 보안 리스크 최소화"),
], size=15.5, gap=12)
# 인증 흐름
fy = Inches(5.2); fh = Inches(0.95)
diagram_node(s, Inches(0.9), fy, Inches(2.2), fh, "사용자", "로그인 요청", fill=GRAY)
arrow(s, Inches(3.2), fy + Inches(0.38), Inches(0.5))
diagram_node(s, Inches(3.8), fy, Inches(2.6), fh, "Cognito Hosted UI", "Google IdP 선택 가능", fill=ACCENT)
arrow(s, Inches(6.5), fy + Inches(0.38), Inches(0.5))
diagram_node(s, Inches(7.1), fy, Inches(2.4), fh, "JWT 발급", "access / refresh", fill=PRIMARY)
arrow(s, Inches(9.6), fy + Inches(0.38), Inches(0.5))
diagram_node(s, Inches(10.2), fy, Inches(2.3), fh, "API 호출", "Bearer 토큰 검증", fill=DARK)

# ============================================================
# 9. 배포 & 인프라
# ============================================================
s = content_slide("04  DEPLOY", "배포 파이프라인 — GitHub Actions → EC2")
fy = Inches(2.3); fh = Inches(1.05)
diagram_node(s, Inches(0.9), fy, Inches(2.3), fh, "git push", "main 브랜치", fill=GRAY)
arrow(s, Inches(3.3), fy + Inches(0.42), Inches(0.5))
diagram_node(s, Inches(3.9), fy, Inches(2.6), fh, "GitHub Actions", "빌드 · 테스트", fill=ACCENT)
arrow(s, Inches(6.6), fy + Inches(0.42), Inches(0.5))
diagram_node(s, Inches(7.2), fy, Inches(2.6), fh, "EC2 배포", "SSH · Docker Compose", fill=PRIMARY)
arrow(s, Inches(9.9), fy + Inches(0.42), Inches(0.5))
diagram_node(s, Inches(10.5), fy, Inches(2.0), fh, "운영 반영", "무중단 재기동", fill=DARK)
bullets(s, Inches(0.9), Inches(4.0), Inches(11.6), [
    ("Monorepo 통합.", "프론트·백엔드·인프라 3개 저장소를 단일 monorepo로 통합해 배포 단일화"),
    ("Docker Compose 운영.", "app + PostgreSQL + Redis 3개 컨테이너 구성, 환경변수로 비밀값 분리"),
    ("push 한 번으로 배포 완료.", "수동 배포 과정 제거 → 시연 직전까지 안전하게 수정 가능"),
], size=15.5, gap=12)

# ============================================================
# 10. 품질 검증
# ============================================================
s = content_slide("05  QUALITY", "품질 검증 — Playwright E2E 72건")
# 숫자 카드
stats = [("72", "전체 시나리오", PRIMARY), ("41", "PASS", GREEN), ("10", "FAIL → 수정", RED), ("21", "SKIP (환경 제약)", GRAY)]
x = Inches(0.9)
for num, label, color in stats:
    add_box(s, x, Inches(2.0), Inches(2.75), Inches(1.5), fill=LIGHT)
    add_text(s, x, Inches(2.2), Inches(2.75), Inches(0.7), [(num, 32, color, True)], align=PP_ALIGN.CENTER)
    add_text(s, x, Inches(3.0), Inches(2.75), Inches(0.4), [(label, 13, GRAY, False)], align=PP_ALIGN.CENTER)
    x += Inches(2.95)
bullets(s, Inches(0.9), Inches(4.1), Inches(11.6), [
    ("인증·홈·채팅·오답노트·친구·설정", "전 영역을 시나리오 기반으로 검증"),
    ("테스트로 발견 → 수정 → 재검증 사이클.", "예: 오답 개수 API 500 에러를 E2E로 발견, 원인(dev 프로필 JWT 타입 불일치) 분석 후 수정"),
    ("회귀 방지.", "수정 후 동일 시나리오 재실행으로 사이드이펙트 없음을 확인"),
], size=15.5, gap=12)

# ============================================================
# 11. 트러블슈팅
# ============================================================
s = content_slide("05  TROUBLESHOOTING", "트러블슈팅 — 증상 · 원인 · 해결")
cases = [
    ("레벨이 항상 그대로인 버그",
     "증상  대화를 잘해도 레벨 변화 없음\n원인  평가 점수 하드코딩(75) + 비동기 평가 결과 미반영 구조\n해결  실제 평가 점수 사용 + 평가 완료 후 레벨 반영 흐름 재설계"),
    ("dev 프로필 NPE (500 에러)",
     "증상  오답노트 API가 500 반환\n원인  @AuthenticationPrincipal Jwt — dev 모드에선 jwt=null\n해결  Long userId로 시그니처 변경, dev/운영 동작 일치"),
    ("배포 시 컨테이너 이름 충돌",
     "증상  재배포마다 docker compose 기동 실패\n원인  이전 컨테이너 잔존으로 이름 충돌\n해결  배포 스크립트에 정리 단계 추가 → 반복 배포 안정화"),
]
x = Inches(0.8)
for t, d in cases:
    add_box(s, x, Inches(2.0), Inches(3.95), Inches(4.4), fill=LIGHT)
    add_text(s, x + Inches(0.25), Inches(2.2), Inches(3.45), Inches(0.8), [(t, 15, PRIMARY, True)])
    tb = add_text(s, x + Inches(0.25), Inches(3.0), Inches(3.45), Inches(3.2), [])
    tf = tb.text_frame
    for j, line in enumerate(d.split("\n")):
        p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
        tag, body = line.split("  ", 1)
        r = p.add_run(); r.text = tag + "  "; set_font(r, 12.5, ACCENT, True)
        r = p.add_run(); r.text = body; set_font(r, 12.5, DARK, False)
        p.space_after = Pt(10)
    x += Inches(4.15)

# ============================================================
# 12. 라이브 시연
# ============================================================
s = content_slide("06  LIVE DEMO", "라이브 시연 — 한 사용자의 학습 흐름")
steps = [
    ("1", "Google 로그인", "Cognito 페더레이션으로 간편 시작"),
    ("2", "홈 대시보드", "통계 카드 · 오늘의 목표 · 오늘의 표현"),
    ("3", "AI 친구와 채팅", "일부러 틀린 문장 → 스트리밍 응답 + 교정 피드백"),
    ("4", "오답노트", "방금 틀린 표현이 자동으로 쌓인 것 확인"),
    ("5", "커스텀 친구 생성", "나만의 AI 대화 상대 만들기"),
    ("6", "홈 복귀", "목표 진행률·통계 자동 갱신 확인"),
]
y = Inches(2.0)
for num, t, d in steps:
    c = add_box(s, Inches(0.9), y, Inches(0.55), Inches(0.55), fill=PRIMARY, shape=MSO_SHAPE.OVAL)
    add_text(s, Inches(0.9), y + Inches(0.1), Inches(0.55), Inches(0.4),
             [(num, 16, WHITE, True)], align=PP_ALIGN.CENTER)
    add_text(s, Inches(1.75), y + Inches(0.03), Inches(3.6), Inches(0.5), [(t, 16, DARK, True)])
    add_text(s, Inches(5.5), y + Inches(0.08), Inches(7.2), Inches(0.5), [(d, 13.5, GRAY, False)])
    y += Inches(0.78)
add_text(s, Inches(0.9), y + Inches(0.1), Inches(11.6), Inches(0.5),
         [[("※ ", 13, RED, True), ("네트워크 장애 대비 시연 녹화 영상 백업 준비", 13, GRAY, False)]])

# ============================================================
# 13. 회고 & 향후 계획
# ============================================================
s = content_slide("07  RETROSPECT", "회고 & 향후 계획")
add_text(s, Inches(0.9), Inches(1.95), Inches(5.6), Inches(0.5), [("배운 점", 18, PRIMARY, True)])
bullets(s, Inches(0.9), Inches(2.5), Inches(5.6), [
    ("", "LLM 프롬프트 설계 — 페르소나·난이도를 프롬프트로 제어하는 경험"),
    ("", "Cognito·Bedrock 등 관리형 서비스로 개발 속도와 보안 양립"),
    ("", "E2E 테스트가 실제 버그를 찾아내는 과정 체득"),
    ("", "monorepo + GitHub Actions로 배포 자동화 구축"),
], size=14, gap=10)
add_text(s, Inches(7.0), Inches(1.95), Inches(5.6), Inches(0.5), [("향후 계획", 18, PRIMARY, True)])
bullets(s, Inches(7.0), Inches(2.5), Inches(5.6), [
    ("", "음성 대화 — STT/TTS로 실제 말하기 연습 지원"),
    ("", "주간 학습 리포트 — 레벨 추이·취약 표현 분석"),
    ("", "오답 기반 복습 퀴즈 자동 생성"),
    ("", "친구(페르소나) 마켓플레이스"),
], size=14, gap=10)
add_box(s, Inches(6.55), Inches(2.0), Pt(2), Inches(4.0), fill=RGBColor(0xE5, 0xE7, 0xEB), shape=MSO_SHAPE.RECTANGLE)

# ============================================================
# 14. Q&A
# ============================================================
s = prs.slides.add_slide(BLANK)
add_box(s, 0, 0, SW, SH, fill=DARK, shape=MSO_SHAPE.RECTANGLE)
add_box(s, Inches(-1.5), Inches(4.8), Inches(5), Inches(5), fill=PRIMARY, shape=MSO_SHAPE.OVAL)
add_box(s, Inches(-0.8), Inches(5.8), Inches(4), Inches(4), fill=ACCENT, shape=MSO_SHAPE.OVAL)
add_text(s, 0, Inches(2.8), SW, Inches(1.2), [("Q & A", 54, WHITE, True)], align=PP_ALIGN.CENTER)
add_text(s, 0, Inches(4.2), SW, Inches(0.6),
         [("감사합니다 — Fluento", 18, RGBColor(0xC7, 0xC5, 0xE8), False)], align=PP_ALIGN.CENTER)

out = "/Users/jeongsua/dev/03_fluento/Fluento_최종발표_초안.pptx"
prs.save(out)
print("saved:", out, "/ slides:", len(prs.slides._sldIdLst))
